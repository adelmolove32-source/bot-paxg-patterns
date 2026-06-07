"""
APRESENTACAO COMPLETA - ESTRATEGIA DE PATTERNS
Com graficos, exemplos e simulado com banca de $200
CORRIGIDO: simulacao realista (retorno diario, nao por trade)
"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

# Cores
PRETO = RGBColor(18, 18, 18)
BRANCO = RGBColor(255, 255, 255)
VERDE = RGBColor(0, 200, 83)
VERMELHO = RGBColor(255, 23, 68)
AZUL = RGBColor(41, 121, 255)
AMARELO = RGBColor(255, 193, 7)
CINZA = RGBColor(158, 158, 158)
ROXO = RGBColor(156, 39, 176)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=PRETO):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=BRANCO, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, 1, 1.5, 11, 1.5, title, size=44, color=VERDE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, 1, 3.2, 11, 1, subtitle, size=24, color=CINZA, align=PP_ALIGN.CENTER)
    return slide

def simulate_balance(initial, pnl_total_pct, days, trades_per_day):
    """Simulacao realista: retorna diario baseado no PnL total do backtest"""
    daily_return = (pnl_total_pct / 100) / days
    balance = [initial]
    for d in range(days):
        new_bal = balance[-1] * (1 + daily_return)
        balance.append(new_bal)
    return balance

# =============================================================================
# SLIDE 1: CAPA
# =============================================================================
slide = add_title_slide(
    "ESTRATEGIA DE PATTERNS",
    "20 Padroes Complexos de Analise Tecnica\nBTC/USDT e ETH/USDT | Binance\n\nBanca: $200 | Backtest: 60 dias"
)

# =============================================================================
# SLIDE 2: O QUE SAO OS PADROES
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "20 PADROES COMPLEXOS", size=36, color=VERDE, bold=True)

patterns_list = [
    ("TRIANGULOS", "ASC_TRI, DESC_TRI, SYM_TRI_BULL, SYM_TRI_BEAR", AZUL),
    ("TOPO/FUNDO", "TRI_TOP, TRI_BOT, DBL_TOP, DBL_BOT", AMARELO),
    ("WEDGE", "RISE_WEDGE, FALL_WEDGE", ROXO),
    ("CABECA E OMBROS", "H&S, INV_H&S", VERMELHO),
    ("BROADENING", "BROAD_TOP, BROAD_BOT", CINZA),
    ("RETANGULO", "BULL_RECT, BEAR_RECT", VERDE),
]

y = 1.4
for cat, items, color in patterns_list:
    add_text(slide, 1, y, 3, 0.5, cat, size=20, color=color, bold=True)
    add_text(slide, 4, y, 8, 0.5, items, size=16, color=CINZA)
    y += 0.6

add_text(slide, 0.5, 5.5, 12, 1, 
    "Cada padrao gera sinais de COMPRA ou VENDA com entry, stop e target definidos.\n"
    "Stop = 50% do range da barra | R:R = 1:3",
    size=16, color=CINZA)

# =============================================================================
# SLIDE 3: EXEMPLO - DESCENDING TRIANGLE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "EXEMPLO: DESCENDING TRIANGLE", size=36, color=VERDE, bold=True)

fig, ax = plt.subplots(figsize=(8, 5))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

x = np.linspace(0, 10, 50)
y = np.sin(x) * 2 + 5 - x * 0.3 + np.random.normal(0, 0.3, 50)

support = np.full(50, 4.5)
resistance = np.linspace(7, 4.5, 50)

ax.plot(x, y, color='#00C853', linewidth=1.5, label='Preco')
ax.plot(x[:35], resistance[:35], color='#FF1744', linewidth=2, linestyle='--', label='Resistencia')
ax.plot(x, support, color='#2979FF', linewidth=2, linestyle='--', label='Suporte')

ax.annotate('ENTRY SHORT', xy=(7.5, 4.6), fontsize=12, color='#FF1744', fontweight='bold')
ax.annotate('STOP', xy=(7.5, 6.5), fontsize=10, color='#FFEB3B')
ax.annotate('TARGET', xy=(7.5, 2.5), fontsize=10, color='#00E676')

ax.axvline(x=7, color='#FFEB3B', linestyle=':', alpha=0.5)
ax.plot(7, 4.5, 'o', color='#FF1744', markersize=12)

ax.set_title('Descending Triangle - Sinal de VENDA', color='white', fontsize=14)
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white')
ax.set_ylabel('Preco', color='white')
ax.set_xlabel('Tempo', color='white')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_desc_tri.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_desc_tri.png', Inches(0.5), Inches(1.3), Inches(7), Inches(5))

add_text(slide, 8, 1.5, 4.5, 1, "COMO FUNCIONA:", size=20, color=VERDE, bold=True)
add_text(slide, 8, 2.2, 4.5, 3,
    "1. Preco forma triangulo descendente\n\n"
    "2. Suporte horizontal + resistencia caindo\n\n"
    "3. BREAKOUT para baixo = SINAL DE VENDA\n\n"
    "4. Entry: na quebra do suporte\n\n"
    "5. Stop: acima do ultimo topo\n\n"
    "6. Target: 3x o risco",
    size=14, color=CINZA)

add_text(slide, 8, 5.5, 4.5, 1, "RESULTADO NO BACKTEST:", size=18, color=AMARELO, bold=True)
add_text(slide, 8, 6.1, 4.5, 0.8,
    "WR: 57.8% | PnL: +205% (60d)\n"
    "Padrao mais lucrativo do 3m",
    size=14, color=VERDE)

# =============================================================================
# SLIDE 4: EXEMPLO - DOUBLE TOP
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "EXEMPLO: DOUBLE TOP", size=36, color=VERDE, bold=True)

fig, ax = plt.subplots(figsize=(8, 5))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

x = np.linspace(0, 10, 60)
y = np.concatenate([
    np.linspace(4, 7, 15),
    np.linspace(7, 5, 10),
    np.linspace(5, 7, 15),
    np.linspace(7, 3.5, 20)
]) + np.random.normal(0, 0.1, 60)

ax.plot(x, y, color='#00C853', linewidth=1.5)
ax.axhline(y=7, color='#FF1744', linewidth=2, linestyle='--', label='Resistencia (Topo 1 e 2)')
ax.axhline(y=5, color='#2979FF', linewidth=2, linestyle='--', label='Suporte (Pescoco)')

ax.annotate('Topo 1', xy=(2, 7.2), fontsize=11, color='#FF1744', fontweight='bold')
ax.annotate('Topo 2', xy=(6.5, 7.2), fontsize=11, color='#FF1744', fontweight='bold')
ax.annotate('ENTRY SHORT', xy=(8, 4.8), fontsize=11, color='#FF1744', fontweight='bold')

ax.plot([2, 6.5], [7, 7], 'o', color='#FF1744', markersize=10)
ax.axvline(x=7.5, color='#FFEB3B', linestyle=':', alpha=0.5)

ax.set_title('Double Top - Padrao de REVERSAO de alta', color='white', fontsize=14)
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white', loc='lower left')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_double_top.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_double_top.png', Inches(0.5), Inches(1.3), Inches(7), Inches(5))

add_text(slide, 8, 1.5, 4.5, 1, "COMO FUNCIONA:", size=20, color=VERDE, bold=True)
add_text(slide, 8, 2.2, 4.5, 3,
    "1. Dois topos na mesma resistencia\n\n"
    "2. Forma letra 'M'\n\n"
    "3. BREAKOUT do pescoco = VENDA\n\n"
    "4. Entry: na quebra do pescoco\n\n"
    "5. Stop: acima do 2o topo\n\n"
    "6. Target: altura do padrao",
    size=14, color=CINZA)

add_text(slide, 8, 5.5, 4.5, 1, "RESULTADO NO BACKTEST:", size=18, color=AMARELO, bold=True)
add_text(slide, 8, 6.1, 4.5, 0.8,
    "WR: 41% | PnL: +81% (60d)\n"
    "Frequente no 3m e 5m",
    size=14, color=VERDE)

# =============================================================================
# SLIDE 5: EXEMPLO - RISE WEDGE
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "EXEMPLO: RISE WEDGE", size=36, color=VERDE, bold=True)

fig, ax = plt.subplots(figsize=(8, 5))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

x = np.linspace(0, 10, 50)
y = 4 + np.sin(x*2) * 0.5 + x * 0.4 + np.random.normal(0, 0.2, 50)

upper = y + np.linspace(1.5, 0.3, 50)
lower = y - np.linspace(0.3, 1.5, 50)

ax.plot(x, y, color='#00C853', linewidth=1.5)
ax.plot(x, upper, color='#FF1744', linewidth=2, linestyle='--', label='Resistencia')
ax.plot(x, lower, color='#2979FF', linewidth=2, linestyle='--', label='Suporte')

ax.annotate('WEDGE\n(Crescente)', xy=(5, 8.5), fontsize=12, color='#9C27B0', fontweight='bold')
ax.annotate('BREAKOUT\nPARA BAIXO', xy=(8.5, 3.5), fontsize=11, color='#FF1744', fontweight='bold')

ax.axvline(x=8, color='#FFEB3B', linestyle=':', alpha=0.5)
ax.plot(8, 5.5, 'o', color='#FF1744', markersize=12)

ax.set_title('Rising Wedge - Padrao de REVERSAO', color='white', fontsize=14)
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_rise_wedge.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_rise_wedge.png', Inches(0.5), Inches(1.3), Inches(7), Inches(5))

add_text(slide, 8, 1.5, 4.5, 1, "COMO FUNCIONA:", size=20, color=VERDE, bold=True)
add_text(slide, 8, 2.2, 4.5, 3,
    "1. Linhas convergem (cada vez mais forca)\n\n"
    "2. Preco sobe mas enfraquece\n\n"
    "3. BREAKOUT para baixo = VENDA\n\n"
    "4. Entry: na quebra do suporte\n\n"
    "5. Stop: acima do ultimo topo\n\n"
    "6. Target: base do wedge",
    size=14, color=CINZA)

add_text(slide, 8, 5.5, 4.5, 1, "RESULTADO NO BACKTEST:", size=18, color=AMARELO, bold=True)
add_text(slide, 8, 6.1, 4.5, 0.8,
    "WR: 60.3% | PnL: +54% (60d)\n"
    "Padrao mais consistente no 30m",
    size=14, color=VERDE)

# =============================================================================
# SLIDE 6: BACKTEST 3m - RANKING
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "BACKTEST 60 DIAS - 3 MINUTOS", size=36, color=VERDE, bold=True)

data_3m = [
    ("Par", "Trades/dia", "WR", "PnL", "PF", "MDD"),
    ("BTC/USDT", "51", "43.6%", "+263%", "1.43", "26.2%"),
    ("ETH/USDT", "49", "43.8%", "+392%", "1.53", "24.1%"),
]

y = 1.5
x_positions = [1, 3.5, 5.5, 7, 8.5, 10]
for i, row in enumerate(data_3m):
    for j, (val, x) in enumerate(zip(row, x_positions)):
        color = VERDE if i == 0 else (AMARELO if j == 3 and i == 2 else BRANCO)
        add_text(slide, x, y, 2, 0.4, val, size=16 if i == 0 else 14, color=color, bold=(i == 0))
    y += 0.5

# Simulacao com $200 - BTC 3m
fig, ax = plt.subplots(figsize=(10, 4))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

days = 60
days_x = np.arange(0, days + 1)

# BTC 3m: +263% em 60 dias = ~4.38% ao dia
btc_3m = simulate_balance(200, 263, days, 51)
# ETH 3m: +392% em 60 dias = ~6.53% ao dia
eth_3m = simulate_balance(200, 392, days, 49)

ax.plot(days_x, btc_3m, color='#F7931A', linewidth=2, label='BTC/USDT 3m')
ax.plot(days_x, eth_3m, color='#627EEA', linewidth=2, label='ETH/USDT 3m')
ax.axhline(y=200, color='#FFEB3B', linestyle='--', alpha=0.5, label='Banca Inicial $200')
ax.fill_between(days_x, 200, btc_3m, alpha=0.15, color='#F7931A')
ax.fill_between(days_x, 200, eth_3m, alpha=0.15, color='#627EEA')

ax.set_title('Evolucao do Saldo - 3 Minutos ($200)', color='white', fontsize=14)
ax.set_xlabel('Dias', color='white')
ax.set_ylabel('Saldo ($)', color='white')
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_equity_3m.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_equity_3m.png', Inches(1.5), Inches(3.5), Inches(10), Inches(4))

add_text(slide, 1.5, 3, 5, 0.5, "SIMULACAO COM $200:", size=16, color=AMARELO, bold=True)
add_text(slide, 6, 3, 6, 0.5, 
    f"BTC: $200 -> ${btc_3m[-1]:,.0f} (+{((btc_3m[-1]/200)-1)*100:.0f}%)  |  "
    f"ETH: $200 -> ${eth_3m[-1]:,.0f} (+{((eth_3m[-1]/200)-1)*100:.0f}%)", 
    size=14, color=VERDE, bold=True)

# =============================================================================
# SLIDE 7: BACKTEST 5m
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "BACKTEST 60 DIAS - 5 MINUTOS", size=36, color=VERDE, bold=True)

data_5m = [
    ("Par", "Trades/dia", "WR", "PnL", "PF", "MDD"),
    ("BTC/USDT", "28", "43.6%", "+167%", "1.38", "27.8%"),
    ("ETH/USDT", "51", "44.5%", "+383%", "1.37", "60.1%"),
]

y = 1.5
for i, row in enumerate(data_5m):
    for j, (val, x) in enumerate(zip(row, x_positions)):
        color = VERDE if i == 0 else BRANCO
        add_text(slide, x, y, 2, 0.4, val, size=14, color=color, bold=(i == 0))
    y += 0.5

fig, ax = plt.subplots(figsize=(10, 4))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

btc_5m = simulate_balance(200, 167, days, 28)
eth_5m = simulate_balance(200, 383, days, 51)

ax.plot(days_x, btc_5m, color='#F7931A', linewidth=2, label='BTC/USDT 5m')
ax.plot(days_x, eth_5m, color='#627EEA', linewidth=2, label='ETH/USDT 5m')
ax.axhline(y=200, color='#FFEB3B', linestyle='--', alpha=0.5, label='Banca Inicial $200')
ax.fill_between(days_x, 200, btc_5m, alpha=0.15, color='#F7931A')
ax.fill_between(days_x, 200, eth_5m, alpha=0.15, color='#627EEA')

ax.set_title('Evolucao do Saldo - 5 Minutos ($200)', color='white', fontsize=14)
ax.set_xlabel('Dias', color='white')
ax.set_ylabel('Saldo ($)', color='white')
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_equity_5m.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_equity_5m.png', Inches(1.5), Inches(3.5), Inches(10), Inches(4))

add_text(slide, 1.5, 3, 5, 0.5, "SIMULACAO COM $200:", size=16, color=AMARELO, bold=True)
add_text(slide, 6, 3, 6, 0.5,
    f"BTC: $200 -> ${btc_5m[-1]:,.0f} (+{((btc_5m[-1]/200)-1)*100:.0f}%)  |  "
    f"ETH: $200 -> ${eth_5m[-1]:,.0f} (+{((eth_5m[-1]/200)-1)*100:.0f}%)",
    size=14, color=VERDE, bold=True)

# =============================================================================
# SLIDE 8: BACKTEST 15m
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "BACKTEST 60 DIAS - 15 MINUTOS", size=36, color=VERDE, bold=True)

data_15m = [
    ("Par", "Trades/dia", "WR", "PnL", "PF", "MDD"),
    ("BTC/USDT", "15", "46.5%", "+331%", "1.85", "38.6%"),
    ("ETH/USDT", "9", "42.9%", "+127%", "1.38", "28.0%"),
]

y = 1.5
for i, row in enumerate(data_15m):
    for j, (val, x) in enumerate(zip(row, x_positions)):
        color = VERDE if i == 0 else BRANCO
        add_text(slide, x, y, 2, 0.4, val, size=14, color=color, bold=(i == 0))
    y += 0.5

fig, ax = plt.subplots(figsize=(10, 4))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

btc_15m = simulate_balance(200, 331, days, 15)
eth_15m = simulate_balance(200, 127, days, 9)

ax.plot(days_x, btc_15m, color='#F7931A', linewidth=2, label='BTC/USDT 15m')
ax.plot(days_x, eth_15m, color='#627EEA', linewidth=2, label='ETH/USDT 15m')
ax.axhline(y=200, color='#FFEB3B', linestyle='--', alpha=0.5, label='Banca Inicial $200')
ax.fill_between(days_x, 200, btc_15m, alpha=0.15, color='#F7931A')
ax.fill_between(days_x, 200, eth_15m, alpha=0.15, color='#627EEA')

ax.set_title('Evolucao do Saldo - 15 Minutos ($200)', color='white', fontsize=14)
ax.set_xlabel('Dias', color='white')
ax.set_ylabel('Saldo ($)', color='white')
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white')
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_equity_15m.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_equity_15m.png', Inches(1.5), Inches(3.5), Inches(10), Inches(4))

add_text(slide, 1.5, 3, 5, 0.5, "SIMULACAO COM $200:", size=16, color=AMARELO, bold=True)
add_text(slide, 6, 3, 6, 0.5,
    f"BTC: $200 -> ${btc_15m[-1]:,.0f} (+{((btc_15m[-1]/200)-1)*100:.0f}%)  |  "
    f"ETH: $200 -> ${eth_15m[-1]:,.0f} (+{((eth_15m[-1]/200)-1)*100:.0f}%)",
    size=14, color=VERDE, bold=True)

# =============================================================================
# SLIDE 9: BACKTEST 30m e 1h
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "BACKTEST 60 DIAS - 30m e 1h", size=36, color=VERDE, bold=True)

data_30m_1h = [
    ("Par", "TF", "Trades/dia", "WR", "PnL", "PF", "MDD"),
    ("BTC", "30m", "8", "44.1%", "+122%", "1.37", "52.5%"),
    ("BTC", "1h", "4", "41.0%", "-36%", "0.84", "55.4%"),
    ("ETH", "30m", "5", "45.4%", "+83%", "1.37", "30.7%"),
    ("ETH", "1h", "3", "50.0%", "+67%", "1.40", "41.8%"),
]

y = 1.5
for i, row in enumerate(data_30m_1h):
    x_pos = [0.5, 2.5, 4.5, 6.5, 8.5, 10.5, 12]
    for j, (val, x) in enumerate(zip(row, x_pos)):
        color = BRANCO
        if i == 0:
            color = VERDE
        elif i == 2 and j == 4:  # -36%
            color = VERMELHO
        add_text(slide, x, y, 2, 0.4, val, size=14, color=color, bold=(i == 0))
    y += 0.5

btc_30m = simulate_balance(200, 122, days, 8)
eth_30m = simulate_balance(200, 83, days, 5)

fig, axes = plt.subplots(1, 2, figsize=(12, 4))
fig.patch.set_facecolor('#121212')

for ax, name, bal in zip(axes, ['BTC 30m', 'ETH 30m'], [btc_30m, eth_30m]):
    ax.set_facecolor('#121212')
    ax.plot(days_x, bal, color='#00C853', linewidth=2)
    ax.axhline(y=200, color='#FFEB3B', linestyle='--', alpha=0.5)
    ax.set_title(f'{name} - $200 -> ${bal[-1]:,.0f}', color='white', fontsize=11)
    ax.tick_params(colors='white')
    for spine in ax.spines.values():
        spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_equity_30m.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_equity_30m.png', Inches(1.5), Inches(4), Inches(10), Inches(3))

add_text(slide, 0.5, 7.2, 12, 0.3, 
    "OBS: BTC 1h ta PERDENDO (-36%) - NAO USAR nesse timeframe",
    size=12, color=VERMELHO)

# =============================================================================
# SLIDE 10: TOP PATTERNS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "TOP PATTERNS POR PERFORMANCE", size=36, color=VERDE, bold=True)

top_patterns = [
    ("DESC_TRI", "Descending Triangle", "57.8%", "+205%", "673", "Mais lucrativo no 3m"),
    ("RISE_WEDGE", "Rising Wedge", "60.3%", "+54%", "73", "Melhor WR no 30m"),
    ("ASC_TRI", "Ascending Triangle", "53.6%", "+153%", "658", "Forte no 3m"),
    ("DBL_TOP", "Double Top", "41.4%", "+81%", "626", "Frequente e consistente"),
    ("SYM_TRI", "Sym. Triangle", "58.4%", "+78%", "315", "Bom WR no 3m"),
]

y = 1.4
headers = ["Pattern", "Nome", "WR", "PnL", "Trades", "Observacao"]
x_pos = [0.5, 2.5, 5, 7, 9, 11]

for j, (h, x) in enumerate(zip(headers, x_pos)):
    add_text(slide, x, y, 2.5, 0.4, h, size=16, color=VERDE, bold=True)
y += 0.5

for row in top_patterns:
    for j, (val, x) in enumerate(zip(row, x_pos)):
        color = VERDE if j == 3 else (AMARELO if j == 2 else BRANCO)
        add_text(slide, x, y, 2.5, 0.4, val, size=13, color=color)
    y += 0.5

add_text(slide, 0.5, 5, 12, 1.5,
    "DESTAQUES:\n"
    "- DESC_TRI: responsavel por ~40% do lucro no 3m\n"
    "- RISE_WEDGE: mais consistente em todos os timeframes\n"
    "- SYM_TRI: bom equilibrio entre WR e PnL",
    size=14, color=CINZA)

# =============================================================================
# SLIDE 11: RANKING FINAL COM $200
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "RANKING FINAL - SIMULACAO $200", size=36, color=VERDE, bold=True)

ranking_data = [
    ("ETH/USDT 3m", 392, 49),
    ("BTC/USDT 15m", 331, 15),
    ("ETH/USDT 5m", 383, 51),
    ("BTC/USDT 3m", 263, 51),
    ("BTC/USDT 5m", 167, 28),
    ("ETH/USDT 30m", 83, 5),
    ("BTC/USDT 30m", 122, 8),
    ("ETH/USDT 1h", 67, 3),
]

ranking_final = []
for name, pnl, tpd in ranking_data:
    final = 200 * (1 + pnl/100)
    ranking_final.append((name, f"${final:,.0f}", f"+{pnl}%", f"{tpd}/dia"))

ranking_final.sort(key=lambda x: float(x[1].replace('$','').replace(',','')), reverse=True)

y = 1.4
headers = ["#", "Timeframe", "Saldo Final", "Retorno", "Trades/dia"]
x_pos = [0.5, 3, 6.5, 9, 11.5]
for h, x in zip(headers, x_pos):
    add_text(slide, x, y, 3, 0.4, h, size=16, color=VERDE, bold=True)
y += 0.5

medals = ["1o", "2o", "3o", "4o", "5o", "6o", "7o", "8o"]
medal_colors = [AMARELO, CINZA, RGBColor(205, 127, 50)] + [BRANCO]*5

for i, (name, final, ret, tpd) in enumerate(ranking_final):
    for val, x, c in zip([medals[i], name, final, ret, tpd], x_pos, 
                          [medal_colors[i]]*5):
        add_text(slide, x, y, 3, 0.4, val, size=13, color=c, bold=(i < 3))
    y += 0.5

# Grafico comparativo
fig, ax = plt.subplots(figsize=(10, 4))
fig.patch.set_facecolor('#121212')
ax.set_facecolor('#121212')

colors = ['#627EEA', '#F7931A', '#627EEA', '#F7931A', '#F7931A', '#627EEA', '#F7931A', '#627EEA']
names = [r[0] for r in ranking_final]
pnls = [392, 331, 383, 263, 167, 83, 122, 67]

for name, pnl, color in zip(names[:5], pnls[:5], colors[:5]):
    bal = simulate_balance(200, pnl, days, 1)
    ax.plot(days_x, bal, linewidth=2, label=name, color=color, alpha=0.8)

ax.axhline(y=200, color='#FFEB3B', linestyle='--', alpha=0.5, label='Banca Inicial')
ax.set_title('Comparacao dos Top 5 - $200 em 60 dias', color='white', fontsize=13)
ax.set_xlabel('Dias', color='white')
ax.set_ylabel('Saldo ($)', color='white')
ax.tick_params(colors='white')
ax.legend(facecolor='#1e1e1e', edgecolor='#333', labelcolor='white', fontsize=9)
for spine in ax.spines.values():
    spine.set_color('#333')

plt.tight_layout()
plt.savefig('chart_ranking.png', dpi=150, bbox_inches='tight', facecolor='#121212')
plt.close()

slide.shapes.add_picture('chart_ranking.png', Inches(1.5), Inches(5.5), Inches(10), Inches(1.8))

# =============================================================================
# SLIDE 12: FILTROS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "FILTROS PARA MELHORAR RESULTADOS", size=36, color=VERDE, bold=True)

filtros = [
    ("FILTRO", "EFEITO", "ONDE FUNCIONA"),
    ("So Bear + Body>50%", "Reduz trades, melhora PF", "ETH 30m/15m"),
    ("So Bull + Body>50%", "MDD baixo", "BTC 15m"),
    ("Body > 60%", "Mantem PnL, reduz ruido", "BTC 30m"),
    ("R:R > 2 ou > 3", "Nao mudou nada", "Nenhum"),
]

y = 1.5
x_f = [0.5, 4, 8]
for i, row in enumerate(filtros):
    for j, (val, x) in enumerate(zip(row, x_f)):
        color = VERDE if i == 0 else BRANCO
        add_text(slide, x, y, 4, 0.4, val, size=14, color=color, bold=(i == 0))
    y += 0.5

add_text(slide, 0.5, 4.5, 12, 1.5,
    "RECOMENDACAO:\n\n"
    "- Usar So Bear + Body>50% para ETH (PF de 1.37 -> 1.49)\n"
    "- Usar So Bull + Body>50% para BTC 15m (MDD de 12% -> 7%)\n"
    "- Filtros de R:R nao funcionam porque o R:R ja e fixo no padrao",
    size=14, color=CINZA)

# =============================================================================
# SLIDE 13: PROXIMOS PASSOS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 12, 0.8, "PROXIMOS PASSOS", size=36, color=VERDE, bold=True)

steps = [
    "1. Configurar API keys no config.py",
    "2. Escolher timeframe (recomendado: 3m ou 15m)",
    "3. Testar com 'python bot.py signals' (ver sinais)",
    "4. Testar com 'python bot.py once' (um ciclo)",
    "5. Rodar em loop com 'python bot.py' (producao)",
    "",
    "RECOMENDACAO FINAL:",
    "- ETH 3m para maximizar lucro ($200 -> $984)",
    "- BTC 15m para menor risco (PF 1.85)",
    "- Comecar com banca pequena ate validar",
]

y = 1.5
for step in steps:
    if step.startswith("REC"):
        add_text(slide, 1, y, 11, 0.5, step, size=18, color=AMARELO, bold=True)
    elif step == "":
        y += 0.2
    else:
        add_text(slide, 1, y, 11, 0.5, step, size=16, color=BRANCO)
    y += 0.5

# =============================================================================
# SLIDE 14: CREDITOS
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 1, 2, 11, 1, "ESTRATEGIA DE PATTERNS", size=44, color=VERDE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11, 1, "20 Padroes Complexos de Analise Tecnica", size=24, color=CINZA, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5, 11, 1, "Dados reais da Binance | Backtest 60 dias", size=18, color=CINZA, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6, 11, 0.5, "Junho 2026", size=16, color=CINZA, align=PP_ALIGN.CENTER)

# =============================================================================
# SALVAR
# =============================================================================
output_path = r'C:\Users\muril\Desktop\bot btc4\estrategia_patterns.pptx'
prs.save(output_path)
print(f"Apresentacao salva em: {output_path}")

# Resumo dos saldos
print("\nSIMULACAO COM $200 EM 60 DIAS:")
for name, final, ret, tpd in ranking_final[:5]:
    print(f"  {name}: $200 -> {final} ({ret})")
